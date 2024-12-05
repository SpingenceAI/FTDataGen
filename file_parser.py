"""Parse data from file"""

import os
import PyPDF2
import docx
import pptx
from bs4 import BeautifulSoup
import markdownify


def extract_text_from_pdf(uploaded_file) -> str:
    """Extract text from PDF file."""
    pdf_reader = PyPDF2.PdfReader(uploaded_file)
    text = ""
    for idx, page in enumerate(pdf_reader.pages):
        if page.extract_text():
            text += f"Page {idx+1}:\n{page.extract_text()}\n\n"
    text = f"## PDF Content\n{text}"
    return text


def extract_text_from_docx(doc_path: str) -> str:
    """
    Parse docx file to markdown format
    """
    content = "## Docx Content\n"
    doc = docx.Document(doc_path)
    element = doc
    table_count = 0
    paragraph_count = 0
    for child in element.element.body:
        if child.tag.endswith("p"):  # Paragraph element
            if child.text == "":
                continue
            content += f"Paragraph({paragraph_count}): {child.text}\n"
            paragraph_count += 1
        elif child.tag.endswith("tbl"):  # Table element
            table = element.tables[table_count]
            table_markdown = ""
            for header in table.rows[0].cells:
                table_markdown += "|" + header.text
            table_markdown += "|\n"
            table_markdown += (
                "|".join(["---" for _ in range(len(table.rows[0].cells))]) + "|\n"
            )
            for row in table.rows[1:]:
                # if all row is empty, skip
                if all(cell.text == "" for cell in row.cells):
                    continue
                row_text = [cell.text for cell in row.cells]
                table_markdown += "|" + "|".join(row_text) + "|\n"
            table_count += 1
            content += f"Table({table_count}): \n{table_markdown}\n"
    return content


def extract_text_from_generic_file(file_path: str) -> str:
    """Extract text from file."""
    with open(file_path, "r") as file:
        return file.read()


def extract_text_from_html(file_path: str) -> str:
    """Extract text from HTML file
    use body tag content to convert to markdown format
    """
    with open(file_path, "r") as file:
        soup = BeautifulSoup(file, "html.parser")
        body = soup.body
        text = markdownify.markdownify(body.decode_contents())
        text = f"## HTML Content\n{text}"
        return text


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    ppt = pptx.Presentation(file_path)
    text = ""
    for idx, slide in enumerate(ppt.slides):
        temp = f"Slide {idx+1}:\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                temp += f"{shape.text}\n\n"
        text += temp
    text = f"## PPTX Content\n{text}"
    return text


def parse_file(file_path: str) -> str:
    """Parse file content to string"""

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"The file {file_path} does not exist")
    if os.path.isdir(file_path):
        raise ValueError(f"The file {file_path} is a directory")

    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension in [".txt", ".md"]:
        return extract_text_from_generic_file(file_path)
    elif file_extension in [".html"]:
        return extract_text_from_html(file_path)
    elif file_extension in [".docx"]:
        return extract_text_from_docx(file_path)
    elif file_extension in [".pptx"]:
        return extract_text_from_pptx(file_path)
    elif file_extension in [".pdf"]:
        return extract_text_from_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {file_extension}")
